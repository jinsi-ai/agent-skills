"""Content: slides, placeholders, bullet points, text extraction."""

import json
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


def _check_pptx():
    if Presentation is None:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")


def add_slide(file_path, layout_index=0, title=None):
    _check_pptx()
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    prs = Presentation(file_path)
    layouts = prs.slide_layouts
    if layout_index < 0 or layout_index >= len(layouts):
        return f"Error: layout_index {layout_index} out of range (0-{len(layouts)-1})"
    slide_layout = layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    if title is not None:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                shape.text = str(title)
                break
    prs.save(file_path)
    return f"Added slide (layout {layout_index})"


def get_slide_info(file_path, slide_index):
    _check_pptx()
    if not os.path.exists(file_path):
        return json.dumps({"error": f"File not found: {file_path}"})
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return json.dumps({"error": f"slide_index {slide_index} out of range (0-{len(slides)-1})"})
    slide = slides[slide_index]
    placeholders = []
    for shape in slide.placeholders:
        placeholders.append({"idx": shape.placeholder_format.idx, "name": shape.name, "text": shape.text})
    shapes = []
    for i, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        shapes.append({"index": i, "name": shape.name, "text": shape.text[:100]})
    info = {
        "slide_index": slide_index,
        "placeholders": placeholders,
        "text_shapes": shapes,
    }
    return json.dumps(info, ensure_ascii=False, indent=2)


def _extract_text_from_slide(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                parts.append(para.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(p for p in parts if p.strip())


def extract_slide_text(file_path, slide_index):
    _check_pptx()
    if not os.path.exists(file_path):
        return json.dumps({"error": f"File not found: {file_path}"})
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return json.dumps({"error": f"slide_index {slide_index} out of range"})
    text = _extract_text_from_slide(slides[slide_index])
    return json.dumps({"slide_index": slide_index, "text": text}, ensure_ascii=False, indent=2)


def extract_presentation_text(file_path, include_slide_info=True):
    _check_pptx()
    if not os.path.exists(file_path):
        return json.dumps({"error": f"File not found: {file_path}"})
    prs = Presentation(file_path)
    slides_text = []
    all_parts = []
    for i, slide in enumerate(prs.slides):
        text = _extract_text_from_slide(slide)
        slides_text.append({"slide_index": i, "text": text})
        if include_slide_info:
            all_parts.append(f"=== Slide {i + 1} ===\n{text}")
        else:
            all_parts.append(text)
    result = {
        "file_path": file_path,
        "total_slides": len(prs.slides),
        "slides": slides_text,
        "all_text_combined": "\n\n".join(all_parts),
    }
    return json.dumps(result, ensure_ascii=False, indent=2)


def populate_placeholder(file_path, slide_index, placeholder_idx, text):
    _check_pptx()
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return f"Error: slide_index {slide_index} out of range"
    slide = slides[slide_index]
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            shape.text = str(text)
            prs.save(file_path)
            return f"Populated placeholder {placeholder_idx}"
    return f"Error: placeholder {placeholder_idx} not found on slide {slide_index}"


def add_bullet_points(file_path, slide_index, items, left=1.0, top=1.5, width=8.0, height=4.0):
    _check_pptx()
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    if not items:
        return "Error: items cannot be empty"
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return f"Error: slide_index {slide_index} out of range"
    slide = slides[slide_index]
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
    prs.save(file_path)
    return f"Added {len(items)} bullet points"
