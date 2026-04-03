"""Presentation management: create, from template, info, core properties, template info."""

import json
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


def _check_pptx():
    if Presentation is None:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")


def _find_template(template_path, search_dirs=None):
    path = Path(template_path)
    if path.exists():
        return str(path)
    if search_dirs is None:
        search_dirs = [".", "./templates", "./assets", "./resources"]
    name = path.name
    for d in search_dirs:
        candidate = Path(d).expanduser().resolve() / name
        if candidate.exists():
            return str(candidate)
    return None


def create_presentation(file_path):
    _check_pptx()
    prs = Presentation()
    prs.save(file_path)
    return f"Created presentation: {file_path}"


def create_presentation_from_template(template_path, file_path, search_dirs=None):
    _check_pptx()
    resolved = _find_template(template_path, search_dirs)
    if not resolved:
        return f"Error: template not found: {template_path}"
    prs = Presentation(resolved)
    prs.save(file_path)
    return f"Created from template: {file_path}"


def get_presentation_info(file_path):
    _check_pptx()
    if not os.path.exists(file_path):
        return json.dumps({"error": f"File not found: {file_path}"})
    prs = Presentation(file_path)
    core = prs.core_properties
    info = {
        "file_path": file_path,
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "title": getattr(core, "title", None) or "",
        "subject": getattr(core, "subject", None) or "",
        "author": getattr(core, "author", None) or "",
        "keywords": getattr(core, "keywords", None) or "",
    }
    return json.dumps(info, ensure_ascii=False, indent=2)


def set_core_properties(file_path, title=None, subject=None, author=None, keywords=None, comments=None):
    _check_pptx()
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    prs = Presentation(file_path)
    core = prs.core_properties
    if title is not None:
        core.title = str(title)
    if subject is not None:
        core.subject = str(subject)
    if author is not None:
        core.author = str(author)
    if keywords is not None:
        core.keywords = str(keywords)
    if comments is not None:
        core.comments = str(comments)
    prs.save(file_path)
    return "Core properties updated"


def get_template_file_info(template_path, search_dirs=None):
    _check_pptx()
    resolved = _find_template(template_path, search_dirs)
    if not resolved:
        return json.dumps({"error": f"Template not found: {template_path}"})
    prs = Presentation(resolved)
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({"idx": ph.placeholder_format.idx, "name": ph.name})
        layouts.append({"index": i, "name": layout.name, "placeholders": placeholders})
    info = {
        "template_path": resolved,
        "slide_count": len(prs.slides),
        "layout_count": len(layouts),
        "layouts": layouts,
    }
    return json.dumps(info, ensure_ascii=False, indent=2)
