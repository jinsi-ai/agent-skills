"""Structural: tables, shapes, charts, table cell formatting."""

import json
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


def _check_pptx():
    if Presentation is None:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")


# Map common shape type names to MSO_AUTO_SHAPE_TYPE (integer) for add_shape
SHAPE_TYPE_MAP = {
    "rectangle": 1,
    "rounded_rectangle": 5,
    "oval": 9,
    "diamond": 4,
    "triangle": 7,
    "arrow": 33,
    "star": 92,
    "flowchart_process": 61,
    "flowchart_decision": 63,
}


def add_table(file_path, slide_index, rows, cols, data=None, left=1.0, top=2.0, width=8.0, height=3.0):
    _check_pptx()
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return f"Error: slide_index {slide_index} out of range"
    slide = slides[slide_index]
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    if data:
        for r, row_data in enumerate(data):
            if r >= rows:
                break
            for c, val in enumerate(row_data):
                if c < cols:
                    table.cell(r, c).text = str(val)
    prs.save(file_path)
    return f"Added table {rows}x{cols}"


def add_shape(file_path, slide_index, shape_type, left=1.0, top=1.0, width=1.0, height=1.0, text=None):
    _check_pptx()
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    shape_type_lower = str(shape_type).lower()
    if shape_type_lower not in SHAPE_TYPE_MAP:
        return f"Error: unsupported shape_type. Use one of: {list(SHAPE_TYPE_MAP.keys())}"
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return f"Error: slide_index {slide_index} out of range"
    slide = slides[slide_index]
    shape_val = SHAPE_TYPE_MAP[shape_type_lower]
    shape = slide.shapes.add_shape(shape_val, Inches(left), Inches(top), Inches(width), Inches(height))
    if text:
        shape.text = str(text)
    prs.save(file_path)
    return f"Added shape: {shape_type}"


def format_table_cell(file_path, slide_index, shape_index, row_index, col_index, text=None, font_size=None, bold=None):
    _check_pptx()
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return f"Error: slide_index {slide_index} out of range"
    slide = slides[slide_index]
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return f"Error: shape_index {shape_index} out of range"
    shape = slide.shapes[shape_index]
    if not shape.has_table:
        return "Error: shape is not a table"
    table = shape.table
    if row_index < 0 or row_index >= table.rows or col_index < 0 or col_index >= table.columns:
        return "Error: row_index or col_index out of range"
    cell = table.cell(row_index, col_index)
    if text is not None:
        cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        if font_size is not None:
            paragraph.font.size = Pt(font_size)
        if bold is not None:
            paragraph.font.bold = bool(bold)
    prs.save(file_path)
    return "Formatted table cell"


def add_chart(file_path, slide_index, chart_type, categories, series_names, series_values, left=1.0, top=2.0, width=8.0, height=4.0, title=None):
    _check_pptx()
    chart_type = str(chart_type).lower()
    if chart_type not in ("column", "bar", "line", "pie"):
        return "Error: chart_type must be column, bar, line, or pie"
    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError:
        return "Error: python-pptx chart support required"
    type_map = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED, "line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}
    prs = Presentation(file_path)
    slides = prs.slides
    if slide_index < 0 or slide_index >= len(slides):
        return f"Error: slide_index {slide_index} out of range"
    slide = slides[slide_index]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in zip(series_names, series_values):
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(type_map[chart_type], Inches(left), Inches(top), Inches(width), Inches(height), chart_data).chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = str(title)
    prs.save(file_path)
    return f"Added {chart_type} chart"
